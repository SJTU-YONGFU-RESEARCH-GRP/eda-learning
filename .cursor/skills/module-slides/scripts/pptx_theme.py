"""Themed PowerPoint rendering for module-slides."""

from __future__ import annotations

import textwrap
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

try:
    from PIL import Image
except ImportError:
    Image = None  # type: ignore[misc, assignment]

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN_L = Inches(0.6)
CONTENT_W = Inches(12.1)
FOOTER_H = Inches(0.35)

FONT_TITLE = Pt(40)
FONT_SUBTITLE = Pt(26)
FONT_HEADING = Pt(28)
FONT_BODY = Pt(20)
FONT_CODE = Pt(20)
FONT_FOOTER = Pt(10)
FONT_SECTION = Pt(36)

FONT_SANS = "Calibri"
FONT_MONO = "Consolas"

LINE_SPACING = Pt(22)
BULLET_SPACE_AFTER = Pt(10)
CONTINUATION_INDENT = "      "

COLOR_TITLE = RGBColor(0x1A, 0x1A, 0x2E)
COLOR_BODY = RGBColor(0x33, 0x33, 0x33)
COLOR_CODE_BG = RGBColor(0xF4, 0xF4, 0xF4)
COLOR_FOOTER = RGBColor(0x66, 0x66, 0x66)

MAX_BULLETS = 6
MAX_CHARS_PER_BULLET = 110
MAX_TITLE_CHARS = 70
# Code / pseudocode slides: hard visual budget (title + monospace + padding).
# Including blank lines — longer blocks must split across slides (no crop).
CODE_SLIDE_MAX_LINES = 12
CODE_SLIDE_WRAP_WIDTH = 68

# Hanging indent for bullet paragraphs.
# Wrap lines must start further right than the bullet glyph + gap, otherwise
# LibreOffice/PowerPoint align wrapped text under the • instead of under body text.
# Asymmetric values: marL = where body text column starts; hanging = bullet zone.
BULLET_LEFT_INDENT = Inches(0.50)  # continuation / wrap margin (marL)
BULLET_HANGING = Inches(-0.25)  # first-line offset so • sits left of body text


def _configure_text_frame(tf: Any) -> None:
    """Disable auto-fit so font sizes stay identical across slides."""
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE


def _set_run_font(
    run: Any,
    size: Pt,
    *,
    font_name: str = FONT_SANS,
    bold: bool = False,
    color: RGBColor | None = None,
) -> None:
    """Apply font settings to a single text run."""
    run.font.size = size
    run.font.name = font_name
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _set_paragraph_font(
    paragraph: Any,
    size: Pt,
    *,
    font_name: str = FONT_SANS,
    bold: bool = False,
    color: RGBColor | None = None,
) -> None:
    paragraph.font.size = size
    paragraph.font.name = font_name
    paragraph.font.bold = bold
    if color is not None:
        paragraph.font.color.rgb = color
    paragraph.line_spacing = LINE_SPACING


def _md_to_plain_and_spans(text: str) -> tuple[str, list[tuple[int, int]]]:
    """Convert ``**bold**`` markdown to plain text plus bold spans.

    Returns:
        plain: text with ``**`` markers removed
        spans: list of ``(start, end)`` character offsets in ``plain`` that are bold
    """
    plain_parts: list[str] = []
    spans: list[tuple[int, int]] = []
    i = 0
    while i < len(text):
        if text.startswith("**", i):
            end = text.find("**", i + 2)
            if end == -1:
                plain_parts.append(text[i:])
                break
            start_plain = sum(len(p) for p in plain_parts)
            content = text[i + 2 : end]
            plain_parts.append(content)
            spans.append((start_plain, start_plain + len(content)))
            i = end + 2
        else:
            plain_parts.append(text[i])
            i += 1
    return "".join(plain_parts), spans


def _strip_md_bold(text: str) -> str:
    """Remove ``**`` markers, keeping inner text."""
    plain, _ = _md_to_plain_and_spans(text)
    return plain


def _runs_for_slice(
    plain: str,
    spans: list[tuple[int, int]],
    start: int,
    end: int,
) -> list[tuple[str, bool]]:
    """Split ``plain[start:end]`` into ``(text, is_bold)`` runs."""
    if start >= end:
        return []
    boundaries = {start, end}
    for s, e in spans:
        if e <= start or s >= end:
            continue
        boundaries.add(max(s, start))
        boundaries.add(min(e, end))
    points = sorted(boundaries)
    runs: list[tuple[str, bool]] = []
    for a, b in zip(points, points[1:], strict=False):
        if a >= b:
            continue
        is_bold = any(s <= a and b <= e for s, e in spans)
        runs.append((plain[a:b], is_bold))
    return runs


def _write_rich_paragraph(
    paragraph: Any,
    text: str,
    size: Pt,
    *,
    font_name: str = FONT_SANS,
    bold: bool = False,
    color: RGBColor | None = None,
    parse_md: bool = True,
) -> None:
    """Write ``text`` into ``paragraph``, converting ``**bold**`` to bold runs."""
    # Remove existing runs so we can rebuild with mixed formatting.
    p_elem = paragraph._element
    for child in list(p_elem):
        if child.tag == qn("a:r"):
            p_elem.remove(child)

    paragraph.line_spacing = LINE_SPACING

    if not parse_md or "**" not in text:
        run = paragraph.add_run()
        run.text = _strip_md_bold(text) if parse_md else text
        _set_run_font(run, size, font_name=font_name, bold=bold, color=color)
        return

    plain, spans = _md_to_plain_and_spans(text)
    if not plain:
        run = paragraph.add_run()
        run.text = ""
        _set_run_font(run, size, font_name=font_name, bold=bold, color=color)
        return

    for segment, is_bold in _runs_for_slice(plain, spans, 0, len(plain)):
        run = paragraph.add_run()
        run.text = segment
        _set_run_font(
            run,
            size,
            font_name=font_name,
            bold=bold or is_bold,
            color=color,
        )


def _write_lines(
    tf: Any,
    lines: list[str],
    size: Pt,
    *,
    font_name: str = FONT_SANS,
    bold: bool = False,
    color: RGBColor | None = None,
    space_after: Pt | None = None,
    parse_md: bool = True,
) -> None:
    _configure_text_frame(tf)
    after = space_after if space_after is not None else Pt(4)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _write_rich_paragraph(
            p,
            line,
            size,
            font_name=font_name,
            bold=bold,
            color=color,
            parse_md=parse_md,
        )
        p.space_after = after


def _truncate(text: str, max_len: int) -> str:
    """Truncate on visible length, preserving ``**bold**`` markers when possible."""
    text = " ".join(text.split())
    plain, spans = _md_to_plain_and_spans(text)
    if len(plain) <= max_len:
        return text
    cut = max_len - 1
    # Rebuild markdown for the truncated plain prefix.
    out: list[str] = []
    i = 0
    while i < cut:
        bold_span = next((s for s in spans if s[0] == i and s[1] <= cut), None)
        if bold_span is not None:
            s, e = bold_span
            out.append(f"**{plain[s:e]}**")
            i = e
            continue
        # If cut lands inside a bold span, close early.
        inside = next((s for s in spans if s[0] < i < s[1]), None)
        if inside is not None:
            out.append(f"**{plain[inside[0]:cut]}**")
            break
        next_bold = min((s[0] for s in spans if s[0] > i), default=cut)
        end = min(cut, next_bold)
        out.append(plain[i:end])
        i = end
    return "".join(out).rstrip() + "…"


def _enable_bullet(paragraph: Any) -> None:
    """Add a visible bullet with a stable font (helps LibreOffice PDF export)."""
    p_pr = paragraph._p.get_or_add_pPr()
    for child in list(p_pr):
        if child.tag in (
            qn("a:buChar"),
            qn("a:buNone"),
            qn("a:buAutoNum"),
            qn("a:buFont"),
            qn("a:buSzPct"),
        ):
            p_pr.remove(child)
    bu_font = OxmlElement("a:buFont")
    bu_font.set("typeface", FONT_SANS)
    bu_sz = OxmlElement("a:buSzPct")
    bu_sz.set("val", "100000")  # 100% of paragraph font size
    bu_char = OxmlElement("a:buChar")
    bu_char.set("char", "\u2022")
    # OOXML order: font → size → char
    p_pr.insert(0, bu_char)
    p_pr.insert(0, bu_sz)
    p_pr.insert(0, bu_font)


def _style_bullet_paragraph(paragraph: Any) -> None:
    """Apply bullet glyph + hanging indent so wraps align under body text.

    Sets ``marL`` / ``indent`` directly on ``a:pPr``. python-pptx property
    setters alone were not persisting indent into the slide XML (so LibreOffice
    fell back to default bullet layout and wraps sat under the •).
    """
    _enable_bullet(paragraph)
    p_pr = paragraph._p.get_or_add_pPr()
    p_pr.set("lvl", "0")
    p_pr.set("marL", str(int(BULLET_LEFT_INDENT)))
    p_pr.set("indent", str(int(BULLET_HANGING)))
    paragraph.space_after = BULLET_SPACE_AFTER


def _add_bullet_paragraphs(tf: Any, bullets: list[str]) -> None:
    """Write one paragraph per bullet; let PowerPoint wrap at the text-box edge."""
    _configure_text_frame(tf)
    for i, raw in enumerate(bullets):
        text = str(raw).strip()
        if not text:
            continue
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _write_rich_paragraph(p, text, FONT_BODY, color=COLOR_BODY)
        _style_bullet_paragraph(p)


def _fit_picture(
    slide: Any,
    img_path: Path,
    left: Any,
    top: Any,
    max_width: Any,
    max_height: Any,
) -> Any:
    """Insert image preserving aspect ratio; center in max box."""
    max_w_in = float(max_width) / float(Inches(1))
    max_h_in = float(max_height) / float(Inches(1))

    if Image is not None:
        with Image.open(img_path) as im:
            iw, ih = im.size
        if iw > 0 and ih > 0:
            ratio = iw / ih
            w_in = max_w_in
            h_in = w_in / ratio
            if h_in > max_h_in:
                h_in = max_h_in
                w_in = h_in * ratio
        else:
            w_in, h_in = max_w_in, max_h_in
    else:
        w_in, h_in = max_w_in, max_h_in

    pic = slide.shapes.add_picture(str(img_path), left, top, width=Inches(w_in))
    pic_h = pic.height
    max_h_emu = int(max_height)
    if pic_h < max_h_emu:
        pic.top = int(top) + (max_h_emu - pic_h) // 2
    return pic


def _add_footer(slide: Any, text: str) -> None:
    if not text:
        return
    box = slide.shapes.add_textbox(
        MARGIN_L, SLIDE_HEIGHT - Inches(0.55), CONTENT_W, FOOTER_H
    )
    tf = box.text_frame
    _configure_text_frame(tf)
    tf.text = text
    p = tf.paragraphs[0]
    _set_paragraph_font(p, FONT_FOOTER, color=COLOR_FOOTER)
    p.alignment = PP_ALIGN.RIGHT


def _add_title_block(
    slide: Any, title: str, subtitle: str | None = None, footer: str = ""
) -> None:
    title = _truncate(title, MAX_TITLE_CHARS)
    box = slide.shapes.add_textbox(MARGIN_L, Inches(2.0), CONTENT_W, Inches(3.0))
    tf = box.text_frame
    _configure_text_frame(tf)
    _write_lines(tf, [title], FONT_TITLE, bold=True, color=COLOR_TITLE)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    if subtitle:
        sub = _truncate(subtitle, 120)
        p2 = tf.add_paragraph()
        _write_rich_paragraph(p2, sub, FONT_SUBTITLE, color=COLOR_BODY)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(14)
    _add_footer(slide, footer)


def _add_heading(slide: Any, title: str, top: float = 0.45) -> None:
    title = _truncate(title, MAX_TITLE_CHARS)
    box = slide.shapes.add_textbox(MARGIN_L, Inches(top), CONTENT_W, Inches(0.85))
    tf = box.text_frame
    _configure_text_frame(tf)
    p = tf.paragraphs[0]
    _write_rich_paragraph(p, title, FONT_HEADING, bold=True, color=COLOR_TITLE)


def _add_bullets(slide: Any, title: str, bullets: list[str], footer: str) -> None:
    _add_heading(slide, title)
    items = [str(b) for b in bullets[:MAX_BULLETS]]

    box = slide.shapes.add_textbox(MARGIN_L, Inches(1.35), CONTENT_W, Inches(5.4))
    tf = box.text_frame
    _add_bullet_paragraphs(tf, items)
    _add_footer(slide, footer)


def _add_image_slide(
    slide: Any,
    media_dir: Path,
    title: str,
    image_rel: str,
    caption: str | None,
    footer: str,
) -> None:
    _add_heading(slide, title)
    img_path = (media_dir / image_rel).resolve()
    if not img_path.is_file():
        img_path = Path(image_rel).resolve()
    if not img_path.is_file():
        raise FileNotFoundError(f"Image not found: {image_rel}")
    _fit_picture(slide, img_path, MARGIN_L, Inches(1.25), CONTENT_W, Inches(5.0))
    if caption:
        cap = slide.shapes.add_textbox(MARGIN_L, Inches(6.45), CONTENT_W, Inches(0.45))
        cap_tf = cap.text_frame
        _configure_text_frame(cap_tf)
        _write_rich_paragraph(
            cap_tf.paragraphs[0],
            _truncate(str(caption), 140),
            Pt(12),
            color=COLOR_BODY,
        )
    _add_footer(slide, footer)


def _looks_like_try_these(code: str) -> bool:
    """Bash try-these: comment-led command groups. Pseudocode listings are denser."""
    lines = [ln for ln in code.replace("\r\n", "\n").split("\n") if ln.strip()]
    if not lines:
        return False
    comment_lines = sum(1 for ln in lines if ln.lstrip().startswith("#"))
    return comment_lines >= 2 and comment_lines >= len(lines) // 3


def _format_try_these_code(code: str) -> list[str]:
    """Format code for a slide.

    Try-these (bash): keep one blank line before each ``#`` comment group.
    Pseudocode / listings: preserve author blanks; do **not** insert a blank
    after every line (that doubled line count and cropped long sketches).
    """
    raw = code.replace("\r\n", "\n").strip("\n").split("\n")
    compact: list[str] = []
    for line in raw:
        if not line.strip():
            if compact and compact[-1] != "":
                compact.append("")
            continue
        compact.append(line.rstrip())
    while compact and compact[-1] == "":
        compact.pop()

    if not _looks_like_try_these(code):
        return compact

    out: list[str] = []
    for i, line in enumerate(compact):
        if (
            line.lstrip().startswith("#")
            and out
            and out[-1] != ""
            and not out[-1].lstrip().startswith("#")
        ):
            out.append("")
        out.append(line)
    while out and out[-1] == "":
        out.pop()
    return out


def _wrap_code_lines(lines: list[str], width: int = CODE_SLIDE_WRAP_WIDTH) -> list[str]:
    """Wrap long lines for monospace display; blanks stay blanks."""
    display: list[str] = []
    for line in lines:
        if not line.strip():
            display.append("")
            continue
        display.extend(textwrap.wrap(line, width=width) or [""])
    return display


def split_code_for_slides(
    code: str,
    *,
    max_lines: int = CODE_SLIDE_MAX_LINES,
) -> list[str]:
    """Split code into chunks that each fit one code slide after wrap.

    Never truncates with ``# ...`` — authors get full text across N slides.
    """
    lines = _format_try_these_code(code)
    if not lines:
        return [""]

    chunks: list[list[str]] = []
    current: list[str] = []

    def flush() -> None:
        nonlocal current
        while current and current[-1] == "":
            current.pop()
        if current:
            chunks.append(current)
        current = []

    for line in lines:
        trial = current + [line]
        if len(_wrap_code_lines(trial)) <= max_lines:
            current.append(line)
            continue
        flush()
        # Single overlong line: hard-wrap into its own chunk(s)
        wrapped = _wrap_code_lines([line])
        if len(wrapped) <= max_lines:
            current = [line]
        else:
            for i in range(0, len(wrapped), max_lines):
                chunks.append(wrapped[i : i + max_lines])
            current = []
    flush()
    return ["\n".join(c) for c in chunks] or [""]


def _add_code_slide(slide: Any, title: str, code: str, footer: str) -> None:
    """Render one code chunk. Caller must pre-split with ``split_code_for_slides``."""
    _add_heading(slide, title)
    lines = _format_try_these_code(code)
    display_lines = _wrap_code_lines(lines)[:CODE_SLIDE_MAX_LINES]

    box = slide.shapes.add_textbox(MARGIN_L, Inches(1.25), CONTENT_W, Inches(5.5))
    tf = box.text_frame
    _write_lines(
        tf,
        display_lines,
        FONT_CODE,
        font_name=FONT_MONO,
        color=COLOR_BODY,
        space_after=Pt(4),
        parse_md=False,
    )
    fill = box.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_CODE_BG
    _add_footer(slide, footer)


def _add_demo_slide(
    slide: Any,
    media_dir: Path,
    title: str,
    command: str,
    screenshot: str | None,
    footer: str,
) -> None:
    _add_heading(slide, title)
    cmd_lines = textwrap.wrap(_truncate(command, 180), width=68) or [""]
    cmd_lines = [f"$ {cmd_lines[0]}"] + [
        f"{CONTINUATION_INDENT}{cont}" for cont in cmd_lines[1:]
    ]

    cmd_box = slide.shapes.add_textbox(MARGIN_L, Inches(1.15), CONTENT_W, Inches(1.0))
    cmd_tf = cmd_box.text_frame
    _write_lines(
        cmd_tf,
        cmd_lines,
        FONT_CODE,
        font_name=FONT_MONO,
        color=COLOR_BODY,
        space_after=Pt(2),
        parse_md=False,
    )
    fill = cmd_box.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_CODE_BG

    img_top = Inches(2.15)
    img_max_h = Inches(4.35)
    if screenshot:
        img_path = (media_dir / screenshot).resolve()
        if img_path.is_file():
            _fit_picture(slide, img_path, MARGIN_L, img_top, CONTENT_W, img_max_h)
        else:
            note = slide.shapes.add_textbox(MARGIN_L, img_top, CONTENT_W, Inches(1))
            note_tf = note.text_frame
            _configure_text_frame(note_tf)
            _write_rich_paragraph(
                note_tf.paragraphs[0],
                f"(screenshot pending: {screenshot})",
                FONT_BODY,
                color=COLOR_BODY,
                parse_md=False,
            )
    _add_footer(slide, footer)


def _add_two_column(
    slide: Any,
    media_dir: Path,
    title: str,
    left: list[str] | str,
    right: str,
    footer: str,
) -> None:
    _add_heading(slide, title)
    left_box = slide.shapes.add_textbox(MARGIN_L, Inches(1.35), Inches(5.8), Inches(5.2))
    tf = left_box.text_frame
    items = left if isinstance(left, list) else [left]
    _add_bullet_paragraphs(tf, [str(x) for x in items[:MAX_BULLETS]])

    img_path = (media_dir / right).resolve()
    if not img_path.is_file():
        img_path = Path(right).resolve()
    if img_path.is_file():
        _fit_picture(slide, img_path, Inches(6.75), Inches(1.35), Inches(5.9), Inches(5.2))
    _add_footer(slide, footer)


def _add_section(slide: Any, title: str, footer: str) -> None:
    title = _truncate(title, MAX_TITLE_CHARS)
    box = slide.shapes.add_textbox(MARGIN_L, Inches(3.0), CONTENT_W, Inches(1.5))
    tf = box.text_frame
    _configure_text_frame(tf)
    p = tf.paragraphs[0]
    _write_rich_paragraph(p, title, FONT_SECTION, bold=True, color=COLOR_TITLE)
    p.alignment = PP_ALIGN.CENTER
    _add_footer(slide, footer)


def build_deck(media_dir: Path, outline: dict[str, Any], output_path: Path) -> int:
    """Build pptx from outline dict; return slide count."""
    footer_default = str(outline.get("footer", ""))

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank = prs.slide_layouts[6]

    for spec in outline.get("slides", []):
        if not isinstance(spec, dict):
            raise ValueError(f"Invalid slide spec: {spec!r}")
        stype = spec.get("type", "bullets")
        footer = str(spec.get("footer", footer_default))
        notes = spec.get("notes")
        slide = prs.slides.add_slide(blank)

        if stype == "title":
            _add_title_block(
                slide,
                str(spec.get("title", outline.get("title", "Untitled"))),
                spec.get("subtitle") or outline.get("subtitle"),
                footer,
            )
        elif stype == "section":
            _add_section(slide, str(spec.get("title", "")), footer)
        elif stype == "bullets":
            _add_bullets(
                slide,
                str(spec.get("title", "")),
                list(spec.get("bullets", [])),
                footer,
            )
        elif stype == "image":
            _add_image_slide(
                slide,
                media_dir,
                str(spec.get("title", "")),
                str(spec["image"]),
                spec.get("caption"),
                footer,
            )
        elif stype == "code":
            base_title = str(spec.get("title", ""))
            chunks = split_code_for_slides(str(spec.get("code", "")))
            # First chunk uses the already-created slide; extras get new slides.
            for ci, chunk in enumerate(chunks):
                if ci == 0:
                    target = slide
                else:
                    target = prs.slides.add_slide(blank)
                part_title = base_title
                if len(chunks) > 1:
                    part_title = f"{base_title} ({ci + 1}/{len(chunks)})"
                _add_code_slide(target, part_title, chunk, footer)
                if notes and ci == 0:
                    target.notes_slide.notes_text_frame.text = str(notes)
                elif ci > 0:
                    target.notes_slide.notes_text_frame.text = ""
            # Skip the shared notes handler below (already applied for chunk 0).
            continue
        elif stype == "demo":
            _add_demo_slide(
                slide,
                media_dir,
                str(spec.get("title", "")),
                str(spec.get("command", "")),
                spec.get("screenshot"),
                footer,
            )
        elif stype == "two_column":
            _add_two_column(
                slide,
                media_dir,
                str(spec.get("title", "")),
                spec.get("left", []),
                str(spec.get("right", "")),
                footer,
            )
        else:
            raise ValueError(f"Unknown slide type: {stype}")

        if notes:
            slide.notes_slide.notes_text_frame.text = str(notes)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return len(prs.slides)
